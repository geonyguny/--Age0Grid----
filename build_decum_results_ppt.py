
#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Build a presentation with results inserted from CSV/JSON artifacts.

Assumptions
- Run from project root (e.g., G:\01_simul)
- Artifacts live under .\outputs and .\outputs\_logs and .\outputs\figs
- Python packages: pandas, python-pptx, matplotlib

Examples
  .\.venv\Scripts\python.exe build_decum_results_ppt.py --out .\outputs --tag ROOM
"""

import argparse
import json
import sys
import re
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt

# ---------- Utilities ----------

def safe_read_csv(p: Path, **kwargs) -> Optional[pd.DataFrame]:
    try:
        if p.exists():
            return pd.read_csv(p, **kwargs)
    except Exception as e:
        print(f"[WARN] Failed to read CSV: {p} ({e})", file=sys.stderr)
    return None

def safe_read_json(p: Path) -> Optional[dict]:
    try:
        if p.exists():
            with p.open("r", encoding="utf-8") as f:
                return json.load(f)
    except Exception as e:
        print(f"[WARN] Failed to read JSON: {p} ({e})", file=sys.stderr)
    return None

def add_title(prs, title_text: str, subtitle: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide.shapes.title.text = title_text
    if subtitle is not None:
        slide.placeholders[1].text = subtitle
    return slide

def add_section(prs, title_text: str, subtitle: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title = slide.shapes.title
    title.text = title_text
    if subtitle:
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(0.6))
        tf = tx.text_frame
        tf.text = subtitle
        for p in tf.paragraphs:
            p.font.size = Pt(16)
    return slide

def df_to_table(slide, df: pd.DataFrame, left=Inches(0.5), top=Inches(1.8), width=Inches(12.5)):
    if df is None or df.empty:
        tx = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        tx.text_frame.text = "표 데이터가 비어 있습니다."
        return

    rows, cols = df.shape
    table = slide.shapes.add_table(rows+1, cols, left, top, width, Inches(0.8+0.2*rows)).table

    # Header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)

    # Body
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            val = df.iat[i, j]
            if isinstance(val, float):
                # Format small decimals compactly
                cell.text = f"{val:.6f}" if abs(val) < 0.1 else f"{val:.4f}"
            else:
                cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)

def add_bullet(slide, lines: List[str], left=Inches(0.6), top=Inches(2.0), width=Inches(12), height=Inches(4.5)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            tf.text = line
            tf.paragraphs[0].font.size = Pt(18)
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 1
            p.font.size = Pt(16)

def add_image(slide, img_path: Path, left=Inches(0.6), top=Inches(2.0), width=Inches(6.5)):
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), left, top, width=width)
    else:
        tx = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        tx.text_frame.text = f"이미지 파일이 없습니다: {img_path.name}"

def human_now_str():
    return datetime.now().strftime("%Y-%m-%d %H:%M")

# ---------- Main builder ----------

def build_prs(outdir: Path, room_tag: str = "ROOM") -> Path:
    outdir.mkdir(parents=True, exist_ok=True)
    logs = outdir / "_logs"
    figs = outdir / "figs"

    prs = Presentation()
    title = f"퇴직연금 인출 시뮬레이션 결과 보고서"
    subtitle = f"단독변수 → 2D → 최적설계 → 행동편향 효과 | 생성: {human_now_str()}"
    add_title(prs, title, subtitle)

    # 0) 파일 존재 점검 슬라이드
    files_needed = [
        outdir / "ROOM_metrics_snapshot.csv",
        logs / "metrics.csv",
        outdir / "OPT_table_core.csv",
        outdir / "BH_table_core.csv",
        outdir / "OPT_BH_compare.csv",
        outdir / "OPT_BH_summary.csv",
        figs / "optimal_points.json",
        outdir / "OPT_BH_LOCK_profile.json",
    ]
    slide = add_section(prs, "파일 점검")
    lines = []
    for p in files_needed:
        lines.append(f"{p}: {'OK' if p.exists() else 'MISSING'}")
    add_bullet(slide, lines)

    # 1) 단독변수 결과
    snap = safe_read_csv(outdir / "ROOM_metrics_snapshot.csv")
    slide = add_section(prs, "1. 단독변수 실험 결과", "상위 지표 요약")
    if snap is not None:
        # 추정 규칙: 'DEV_1D_' 혹은 'SINGLE_' 등이 태그에 존재
        one_d = snap[snap["tag"].astype(str).str.contains(r"(DEV_1D_|SINGLE_)", regex=True, na=False)] if "tag" in snap.columns else None
        if one_d is not None and not one_d.empty:
            # 핵심 열만
            keep = [c for c in ["tag","EW","ES95","Ruin","mean_WT","w_max","h_fx","ann_alpha"] if c in one_d.columns]
            df = one_d[keep].copy().sort_values(["EW","ES95"], ascending=[False, True]).head(12)
            df_to_table(slide, df)
        else:
            add_bullet(slide, ["단독변수 표식(DEV_1D_|SINGLE_)을 찾지 못했습니다. 스냅샷 열/태그를 확인하세요."])
    else:
        add_bullet(slide, ["ROOM_metrics_snapshot.csv 읽기 실패"])

    # 2) 2D 결과
    slide = add_section(prs, "2. 2D 변수 상호작용 결과", "상위 지표 요약")
    if snap is not None:
        two_d = snap[snap["tag"].astype(str).str.contains(r"(DEV_2D_|2D_)", regex=True, na=False)] if "tag" in snap.columns else None
        if two_d is not None and not two_d.empty:
            keep = [c for c in ["tag","EW","ES95","Ruin","mean_WT","w_max","h_fx","ann_alpha"] if c in two_d.columns]
            df = two_d[keep].copy().sort_values(["EW","ES95"], ascending=[False, True]).head(12)
            df_to_table(slide, df)
        else:
            add_bullet(slide, ["2D 표식(DEV_2D_|2D_)을 찾지 못했습니다. 스냅샷 열/태그를 확인하세요."])
    else:
        add_bullet(slide, ["ROOM_metrics_snapshot.csv 읽기 실패"])

    # 3) 최적설계 요약
    slide = add_section(prs, "3. 최적설계 요약 (OPT)")
    opt_core = safe_read_csv(outdir / "OPT_table_core.csv")
    df_to_table(slide, opt_core)

    # 3-1) 최적설계 파라미터(가능 시 optimal_points.json)
    slide = add_section(prs, "최적 포인트 파라미터 (optional)")
    opt_json = safe_read_json(figs / "optimal_points.json")
    if opt_json:
        pretty = json.dumps(opt_json, ensure_ascii=False, indent=2)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.5), Inches(5.5))
        tf = tx.text_frame
        tf.text = pretty
        tf.paragraphs[0].font.size = Pt(12)
    else:
        add_bullet(slide, ["optimal_points.json 미존재 또는 파싱 실패"])

    # 4) 행동편향 효과 (BH)
    slide = add_section(prs, "4. 최적설계 대비 행동편향 효과 (BH)")
    bh_core = safe_read_csv(outdir / "BH_table_core.csv")
    df_to_table(slide, bh_core)

    # 4-1) OPT vs BH 비교
    slide = add_section(prs, "OPT vs BH 비교 (Top Picks)")
    compare = safe_read_csv(outdir / "OPT_BH_compare.csv")
    df_to_table(slide, compare)

    # 5) 결론/요약 (자동 생성 텍스트)
    slide = add_section(prs, "5. 결론 및 시사점")
    bullets = []
    if opt_core is not None and not opt_core.empty:
        best_opt = opt_core.sort_values(["EW","ES95"], ascending=[False, True]).head(1)
        tag = best_opt.iloc[0]["tag"]
        ew  = best_opt.iloc[0].get("EW", None)
        es  = best_opt.iloc[0].get("ES95", None)
        bullets.append(f"최적설계 기준 태그: {tag}")
        if ew is not None and es is not None:
            bullets.append(f"핵심지표: EW={ew:.6f}, ES95={es:.6f}")
    else:
        bullets.append("최적설계 표가 비어 있어 결론 자동작성에 제약이 있습니다.")

    if bh_core is not None and not bh_core.empty:
        bullets.append("행동편향 적용 후 소비안정지표(la_sf_mean/la_sf_rate)를 함께 비교하여 보수성/안정성 변화를 해석.")
    else:
        bullets.append("BH 핵심표가 비어 있어 효과 비교 슬라이드는 골격만 포함됩니다.")

    bullets.append("ES95(↓ 개선), Ruin(↓ 개선), EW(↑ 개선) 순서로 정책 우열을 종합 평가.")
    bullets.append("1D/2D에서 확인된 우수 영역은 행동편향 적용 후에도 골격 유지. 다만 파라미터 조정 폭은 달라질 수 있음.")
    add_bullet(slide, bullets)

    # 6) (선택) 존재하는 이미지 삽입: figs 폴더의 png
    pngs = sorted(figs.glob("*.png"))
    if pngs:
        for p in pngs[:6]:  # 과도한 슬라이드 증가 방지
            s = add_section(prs, f"부록: {p.name}")
            add_image(s, p, width=Inches(10))

    # Save
    stamp = datetime.now().strftime("%Y%m%d_%H%M")
    out_ppt = outdir / f"Paper_Decum_Report_{stamp}.pptx"
    prs.save(str(out_ppt))
    print(f"[OK] Exported: {out_ppt}")
    return out_ppt

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="./outputs", help="outputs directory")
    ap.add_argument("--tag", default="ROOM", help="room tag prefix (for filtering snapshots, optional)")
    args = ap.parse_args()

    outdir = Path(args.out).resolve()
    out_ppt = build_prs(outdir, room_tag=args.tag)
    print(out_ppt)

if __name__ == "__main__":
    main()
