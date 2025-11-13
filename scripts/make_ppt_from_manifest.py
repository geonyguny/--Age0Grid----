from __future__ import annotations
import argparse, os, pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

def add_title(prs, title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    if subtitle:
        s.placeholders[1].text = subtitle

def add_item_slide(prs, title, img_path: str, table_path: str | None):
    s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    s.shapes.title.text = title

    y = Inches(1.2)
    if img_path and os.path.exists(img_path):
        # 가로폭 10인치 기준, 여백 포함 9.0인치 정도로 배치
        s.shapes.add_picture(img_path, Inches(0.5), y, width=Inches(9.0))
        y = Inches(5.6)

    if table_path and os.path.exists(table_path):
        box = s.shapes.add_textbox(Inches(0.5), y, Inches(9.0), Inches(1.2))
        tf = box.text_frame
        tf.text = f"Table CSV: {table_path}"
        tf.paragraphs[0].font.size = Pt(12)

def main():
    ap = argparse.ArgumentParser(description="Build PPT from MANIFEST_plots_tables.csv")
    ap.add_argument("--manifest", default="MANIFEST_plots_tables.csv")
    ap.add_argument("--title", default="1D/2D Results Summary")
    ap.add_argument("--outfile", default="DEV_1D_Report_from_manifest.pptx")
    args = ap.parse_args()

    if not os.path.exists(args.manifest):
        raise SystemExit(f"[ERR] manifest not found: {args.manifest}")

    df = pd.read_csv(args.manifest)
    # 기대 컬럼: title, plot, table (plot/table 경로는 plot_and_tables.py가 채움)
    for c in ["title","plot","table"]:
        if c not in df.columns:
            raise SystemExit(f"[ERR] manifest missing column: {c}")

    prs = Presentation()
    add_title(prs, args.title, "Auto-generated from plots & tables manifest")

    for _, r in df.iterrows():
        title = str(r.get("title","")).strip() or "(no title)"
        plot  = str(r.get("plot","")).strip()
        table = str(r.get("table","")).strip()
        add_item_slide(prs, title, plot if os.path.exists(plot) else None,
                       table if os.path.exists(table) else None)

    prs.save(args.outfile)
    print(f"[OK] saved: {os.path.abspath(args.outfile)}")

if __name__ == "__main__":
    main()
